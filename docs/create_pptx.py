"""Generate NANOG PowerPoint presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor as RgbColor

# Create presentation (16:9 aspect ratio)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BLUE = RgbColor(0, 102, 204)
GREEN = RgbColor(76, 175, 80)
YELLOW = RgbColor(255, 152, 0)
RED = RgbColor(244, 67, 54)
DARK = RgbColor(30, 30, 30)
GRAY = RgbColor(102, 102, 102)

def add_title_slide(title, subtitle="", author="", event=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.333), Inches(0.8))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(28)
        p2.font.color.rgb = GRAY
        p2.alignment = PP_ALIGN.CENTER

    if author:
        txBox3 = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.333), Inches(0.5))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = author
        p3.font.size = Pt(24)
        p3.alignment = PP_ALIGN.CENTER

    if event:
        txBox4 = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12.333), Inches(0.5))
        tf4 = txBox4.text_frame
        p4 = tf4.paragraphs[0]
        p4.text = event
        p4.font.size = Pt(18)
        p4.font.color.rgb = GRAY
        p4.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(title, bullets=None, code=None, note=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title with underline
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(51, 51, 51)

    # Line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    y_pos = 1.4

    if bullets:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12.333), Inches(4))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf2.paragraphs[0]
            else:
                p = tf2.add_paragraph()
            p.text = "• " + bullet
            p.font.size = Pt(24)
            p.space_after = Pt(12)
        y_pos += len(bullets) * 0.6

    if code:
        # Dark code box
        code_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(y_pos), Inches(12.333), Inches(2.5)
        )
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = DARK
        code_box.line.fill.background()

        txBox3 = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.2), Inches(12), Inches(2.3))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = code
        p3.font.size = Pt(14)
        p3.font.name = "Courier New"
        p3.font.color.rgb = RgbColor(212, 212, 212)
        y_pos += 2.7

    if note:
        # Highlight box
        note_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(y_pos), Inches(12.333), Inches(0.8)
        )
        note_box.fill.solid()
        note_box.fill.fore_color.rgb = RgbColor(255, 243, 224)
        note_box.line.color.rgb = YELLOW

        txBox4 = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.15), Inches(12), Inches(0.6))
        tf4 = txBox4.text_frame
        p4 = tf4.paragraphs[0]
        p4.text = note
        p4.font.size = Pt(18)

    return slide

def add_table_slide(title, headers, rows):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(12.333), Pt(4))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()

    # Table
    cols = len(headers)
    table = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5 * (len(rows) + 1))).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        para = cell.text_frame.paragraphs[0]
        para.font.color.rgb = RgbColor(255, 255, 255)
        para.font.bold = True
        para.font.size = Pt(16)

    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(14)
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(249, 249, 249)

    return slide

# ============================================================
# CREATE SLIDES
# ============================================================

# Slide 1: Title
add_title_slide(
    "Peer Risk Intelligence",
    "Should You Peer With That ASN?",
    "[Your Name]",
    "NANOG [XX] | [Date]"
)

# Slide 2: The Problem
add_content_slide(
    "The Problem",
    bullets=[
        '"Is this network safe to peer with?"',
        "Peering decisions based on PeeringDB + gut feel",
        "No visibility into routing behavior before peering",
        "Discover problems after the BGP session is up",
        "Route leaks and hijacks from peers are expensive"
    ]
)

# Slide 3: Real Incident
add_content_slide(
    "Real Example: Cloudflare 1.1.1.1 (June 2024)",
    bullets=[
        "AS267613 (Eletronet) announced 1.1.1.1/32",
        "AS262504 leaked 1.1.1.0/24 to upstreams",
        "Affected traffic for 7+ hours",
        "Would you have peered with AS267613?"
    ],
    note="Question: Could we have known this network was risky before the incident?"
)

# Slide 4: Introducing Route Sherlock
add_content_slide(
    "Introducing Route Sherlock",
    bullets=[
        "Open-source CLI tool for Peer Risk Intelligence",
        "Data: RIPEstat + PeeringDB + BGPStream + Claude AI",
        "Features: Risk scoring, backtesting, IX overlap, AI reports"
    ]
)

# Slide 5: The Command
add_content_slide(
    "The Command",
    code="""$ route-sherlock peer-risk AS64500

# Options:
$ route-sherlock peer-risk AS64500 --my-asn AS13335  # IX overlap
$ route-sherlock peer-risk AS64500 --days 180        # Extended history
$ route-sherlock peer-risk AS64500 --ai              # AI analysis"""
)

# Slide 6: Demo Cloudflare
add_content_slide(
    "Demo: Scoring Cloudflare (AS13335)",
    code="""$ route-sherlock peer-risk AS13335

+==================== Peer Risk Score ====================+
| 100/100 (100.0%)                                        |
| Risk Level: LOW                                         |
| Recommendation: RECOMMENDED                             |
+=========================================================+

Network: CLOUDFLARENET | Policy: Open | IXes: 350""",
    note="Perfect score - safe to peer"
)

# Slide 7: Demo Eletronet
add_content_slide(
    "Demo: Scoring Eletronet (AS267613)",
    code="""$ route-sherlock peer-risk AS267613

+==================== Peer Risk Score ====================+
| 72/100 (72.0%)                                          |
| Risk Level: MODERATE                                    |
| Recommendation: ACCEPTABLE WITH MONITORING              |
+=========================================================+

Stability: 5/30 - High churn: 1637 updates/day (-25)
Warning: High BGP churn detected""",
    note="Flagged BEFORE the incident occurred!"
)

# Slide 8: Scoring Algorithm
add_table_slide(
    "Scoring Algorithm (100 points)",
    ["Category", "Max", "What We Check"],
    [
        ["Maturity", "20", "PeeringDB presence, IRR as-set, policy URL, IX count"],
        ["Stability", "30", "BGP update frequency (churn detection)"],
        ["Incident History", "30", "Upstream diversity, topology redundancy"],
        ["Policy", "10", "Open (+10), Selective (+7), Restrictive (+3)"],
        ["Security", "10", "IRR registration, transit relationships"],
    ]
)

# Slide 9: Risk Levels
add_table_slide(
    "Risk Levels & Recommendations",
    ["Score", "Risk Level", "Recommendation"],
    [
        ["80-100", "LOW", "Recommended - standard peering process"],
        ["60-79", "MODERATE", "Acceptable - implement monitoring"],
        ["40-59", "ELEVATED", "Caution - strict prefix limits, IRR filtering"],
        ["0-39", "HIGH", "Not recommended - decline or require remediation"],
    ]
)

# Slide 10: Stability Deep Dive
add_content_slide(
    "Stability Score: BGP Churn Detection",
    bullets=[
        "Query RIPEstat for BGP updates over time",
        "Calculate updates per day average",
        "Thresholds: <10/day stable, >100/day high churn",
        "AS267613 had 1,637 updates/day = 5/30 score"
    ]
)

# Slide 11: IX Overlap
add_content_slide(
    "IX Overlap Analysis",
    code="""$ route-sherlock peer-risk AS15169 --my-asn AS13335

## IX Overlap
   Common IXes: 164
   Your IXes: 350 | Target IXes: 198
   Can peer at 164 location(s)""",
    note="Instantly see where you can peer - useful for peering coordinators"
)

# Slide 12: Backtesting
add_content_slide(
    "Historical Backtesting",
    code="""$ route-sherlock backtest 1.1.1.0/24 --origin AS13335 \\
    --time "2024-06-27 18:00" --duration 8h

Anomalies Detected: 329

#1 [HIGH] ROUTE LEAK
   Time: 2024-06-27T18:49:06
   AS Path: 50763 > 1031 > 262504 > 267613 > 13335"""
)

# Slide 13: Validation
add_content_slide(
    "Validation: Did It Work?",
    bullets=[
        "Cloudflare reported: Incident started ~18:51 UTC",
        "Route Sherlock found: First anomaly 18:49:06 UTC",
        "Correctly identified AS267613 and AS262504",
        "Duration match: 7.6 hours vs 7.5 hours reported"
    ],
    note="Detected 2 minutes BEFORE Cloudflare's reported start time"
)

# Slide 14: AI Synthesis
add_content_slide(
    "AI-Powered Analysis (Optional)",
    code="""$ route-sherlock peer-risk AS267613 --ai

+--------------- AI-Generated Risk Assessment ---------------+
| Executive Summary:                                         |
| Conditional peer. AS267613 shows concerning stability      |
| metrics with 1,637 BGP updates/day.                        |
|                                                            |
| Technical Safeguards:                                      |
| - Max-prefix limit: 15 (they announce 7 prefixes)          |
| - Require IRR filtering against RADB::AS-267613            |
| - Enable RPKI-invalid rejection                            |
+------------------------------------------------------------+"""
)

# Slide 15: Comparison
add_table_slide(
    "How Is This Different?",
    ["Feature", "BGPalerter", "ARTEMIS", "Route Sherlock"],
    [
        ["Real-time monitoring", "Yes", "Yes", "No"],
        ["Historical backtesting", "No", "No", "Yes"],
        ["Peer risk scoring", "No", "No", "Yes"],
        ["AI analysis", "No", "No", "Yes"],
        ['"Should I peer?"', "No", "No", "Yes"],
        ["Open source", "Yes", "Yes", "Yes"],
    ]
)

# Slide 16: Practical Safeguards
add_table_slide(
    "Practical Safeguards by Risk Level",
    ["Risk", "Max-Prefix", "IRR Filter", "RPKI"],
    [
        ["LOW", "2x announced", "Standard", "Warn on invalid"],
        ["MODERATE", "1.5x announced", "Strict", "Reject invalid"],
        ["ELEVATED", "1.2x announced", "Strict + verify", "Reject invalid"],
        ["HIGH", "Decline", "N/A", "N/A"],
    ]
)

# Slide 17: Getting Started
add_content_slide(
    "Getting Started",
    code="""# Install
$ pip install route-sherlock

# For historical backtesting
$ brew install bgpstream
$ pip install pybgpstream

# Optional: AI synthesis
$ export ANTHROPIC_API_KEY="your-key"

# Run
$ route-sherlock peer-risk AS64500"""
)

# Slide 18: Summary
add_content_slide(
    "Summary",
    bullets=[
        "Problem: No tool answers 'Should I peer with this ASN?'",
        "Solution: Route Sherlock - Peer Risk Intelligence",
        "How: RIPEstat + PeeringDB + BGPStream + AI",
        "Validated: Correctly flags networks involved in real incidents",
        "Practical: Outputs actionable recommendations"
    ],
    note="Try it: Score a network you're considering peering with today"
)

# Slide 19: Q&A
add_title_slide(
    "Questions?",
    "GitHub: [your-repo-url]",
    "[your-email]",
    "[your-twitter/linkedin]"
)

# Slide 20: Backup Architecture
add_content_slide(
    "Backup: Architecture",
    code="""route_sherlock/
├── cli/
│   ├── main.py          # Typer CLI entry point
│   └── commands.py      # Command implementations
├── collectors/
│   ├── ripestat.py      # RIPEstat API client
│   ├── peeringdb.py     # PeeringDB API client
│   └── bgpstream.py     # Historical BGP archives
└── synthesis/
    └── engine.py        # AI synthesis with Claude""",
    note="Python 3.11+ | Rich terminal UI | pybgpstream for RouteViews/RIPE RIS"
)

# Save
prs.save('/Users/jd/Documents/route-sherlock/docs/NANOG-Slides.pptx')
print("Created: NANOG-Slides.pptx")
